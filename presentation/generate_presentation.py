"""Generate and validate the Spanish executive presentation."""

from __future__ import annotations

import json
from collections import defaultdict
from io import BytesIO
from itertools import pairwise
from pathlib import Path
from typing import Any

from PIL import Image, ImageOps
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.presentation import Presentation as PresentationType
from pptx.slide import Slide
from pptx.util import Inches, Pt

from renewable_operations.data_quality import assert_quality, validate_dataset
from renewable_operations.synthetic_data import generate_dataset
from renewable_operations.transformations import build_daily_kpis

ROOT = Path(__file__).resolve().parents[1]
OUTPUT = ROOT / "presentation" / "renewable_operations_demo.pptx"
LOCAL_EVIDENCE = ROOT / "evidence" / "local_presentation_metrics.json"
REMOTE_EVIDENCE = ROOT / "evidence" / "remote_presentation_metrics.json"
ASSETS = ROOT / "presentation" / "assets"
DASHBOARD_SCREENSHOT = ASSETS / "dashboard_executive_v2.png"
DASHBOARD_RELIABILITY_SCREENSHOT = ASSETS / "dashboard_reliability_v2.png"
GENIE_CONVERSATION_SCREENSHOT = ASSETS / "genie_conversation.png"
GENIE_SCREENSHOT = GENIE_CONVERSATION_SCREENSHOT
GENIE_ONE_HOME_SCREENSHOT = ASSETS / "genie_one_home.png"
GENIE_BENCHMARK_SCREENSHOT = ASSETS / "genie_benchmark_results.png"

NAVY = RGBColor(0x1B, 0x25, 0x33)
RED = RGBColor(0xFF, 0x36, 0x21)
TEAL = RGBColor(0x00, 0xA9, 0x72)
GRAY = RGBColor(0xF5, 0xF7, 0xF9)
MID_GRAY = RGBColor(0x6B, 0x77, 0x85)
LIGHT_LINE = RGBColor(0xD9, 0xDF, 0xE5)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ORANGE = RGBColor(0xF3, 0x9C, 0x4A)

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)
FOOTER = "Datos exclusivamente sintéticos · Julio 2026"


def _textbox(
    slide: Slide,
    text: str,
    x: float,
    y: float,
    width: float,
    height: float,
    *,
    size: float = 18,
    color: RGBColor = NAVY,
    bold: bool = False,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    valign: MSO_ANCHOR = MSO_ANCHOR.MIDDLE,
) -> Any:
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(width), Inches(height))
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = valign
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    run.font.name = "Aptos"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return shape


def _box(
    slide: Slide,
    x: float,
    y: float,
    width: float,
    height: float,
    *,
    fill: RGBColor = WHITE,
    line: RGBColor = LIGHT_LINE,
    radius: bool = True,
) -> Any:
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(1)
    return shape


def _title(slide: Slide, title: str, subtitle: str | None = None) -> None:
    _textbox(slide, title, 0.65, 0.32, 11.9, 0.55, size=26, bold=True)
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.65), Inches(0.98), Inches(1.1), Inches(0.06)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = RED
    line.line.fill.background()
    if subtitle:
        _textbox(slide, subtitle, 1.95, 0.86, 10.5, 0.3, size=11, color=MID_GRAY)


def _footer(slide: Slide, number: int) -> None:
    _textbox(slide, FOOTER, 0.65, 7.12, 6.5, 0.2, size=9, color=MID_GRAY)
    _textbox(
        slide, str(number), 12.1, 7.12, 0.55, 0.2, size=9, color=MID_GRAY, align=PP_ALIGN.RIGHT
    )


def _base_slide(presentation: PresentationType, title: str, subtitle: str | None = None) -> Slide:
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    _title(slide, title, subtitle)
    _footer(slide, len(presentation.slides))
    return slide


def _notes(slide: Slide, message: str, transition: str, duration: str, demo: str = "N/A") -> None:
    slide.notes_slide.notes_text_frame.text = (
        f"Mensaje principal: {message}\n"
        f"Transición: {transition}\n"
        f"Duración aproximada: {duration}\n"
        f"Demostración: {demo}"
    )


def _card(
    slide: Slide,
    title: str,
    body: str,
    x: float,
    y: float,
    width: float,
    height: float,
    *,
    accent: RGBColor = RED,
) -> None:
    _box(slide, x, y, width, height, fill=GRAY)
    accent_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(0.08), Inches(height)
    )
    accent_shape.fill.solid()
    accent_shape.fill.fore_color.rgb = accent
    accent_shape.line.fill.background()
    _textbox(slide, title, x + 0.22, y + 0.12, width - 0.35, 0.38, size=15, bold=True)
    _textbox(
        slide,
        body,
        x + 0.22,
        y + 0.55,
        width - 0.35,
        height - 0.68,
        size=12,
        color=MID_GRAY,
        valign=MSO_ANCHOR.TOP,
    )


def _dark_slide(presentation: PresentationType, title: str, subtitle: str | None = None) -> Slide:
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = NAVY
    _textbox(slide, title, 0.72, 0.42, 11.9, 0.62, size=27, color=WHITE, bold=True)
    red_line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.72), Inches(1.13), Inches(1.15), Inches(0.06)
    )
    red_line.fill.solid()
    red_line.fill.fore_color.rgb = RED
    red_line.line.fill.background()
    if subtitle:
        _textbox(slide, subtitle, 2.05, 1.02, 10.45, 0.32, size=11, color=LIGHT_LINE)
    _textbox(slide, FOOTER, 0.72, 7.12, 8.0, 0.2, size=9, color=LIGHT_LINE)
    _textbox(
        slide,
        str(len(presentation.slides)),
        12.08,
        7.12,
        0.55,
        0.2,
        size=9,
        color=LIGHT_LINE,
        align=PP_ALIGN.RIGHT,
    )
    return slide


def _pill(
    slide: Slide,
    text: str,
    x: float,
    y: float,
    width: float,
    *,
    fill: RGBColor = NAVY,
    color: RGBColor = WHITE,
) -> None:
    _box(slide, x, y, width, 0.4, fill=fill, line=fill)
    _textbox(
        slide,
        text,
        x + 0.08,
        y + 0.02,
        width - 0.16,
        0.3,
        size=10,
        color=color,
        bold=True,
        align=PP_ALIGN.CENTER,
    )


def _kpi_card(
    slide: Slide,
    value: str,
    label: str,
    x: float,
    y: float,
    width: float,
    *,
    accent: RGBColor = TEAL,
) -> None:
    _box(slide, x, y, width, 1.25, fill=GRAY)
    _textbox(slide, value, x + 0.16, y + 0.16, width - 0.32, 0.55, size=24, bold=True)
    _textbox(slide, label, x + 0.16, y + 0.77, width - 0.32, 0.28, size=10, color=MID_GRAY)
    marker = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(0.07), Inches(1.25)
    )
    marker.fill.solid()
    marker.fill.fore_color.rgb = accent
    marker.line.fill.background()


def _source_note(slide: Slide, text: str) -> None:
    _textbox(
        slide,
        text,
        0.78,
        6.83,
        10.9,
        0.18,
        size=8,
        color=MID_GRAY,
        align=PP_ALIGN.RIGHT,
    )


def _arrow(
    slide: Slide,
    start: tuple[float, float],
    end: tuple[float, float],
    *,
    color: RGBColor = LIGHT_LINE,
    width: float = 1.8,
) -> None:
    connector = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(start[0]),
        Inches(start[1]),
        Inches(end[0]),
        Inches(end[1]),
    )
    connector.line.color.rgb = color
    connector.line.width = Pt(width)


def _real_screenshot(
    slide: Slide,
    path: Path,
    x: float,
    y: float,
    width: float,
    height: float,
    *,
    crop_top_pixels: int = 0,
) -> None:
    """Insert a sanitized application capture without distorting its aspect ratio."""
    if not path.exists():
        raise FileNotFoundError(f"Required screenshot not found: {path}")

    with Image.open(path) as source:
        if source.width < 800 or source.height < 600:
            raise ValueError(f"Screenshot dimensions are too small: {path}")
        crop_top = min(crop_top_pixels, source.height - 1)
        cropped = source.crop((0, crop_top, source.width, source.height))
        target_size = (max(1, round(width * 180)), max(1, round(height * 180)))
        fitted = ImageOps.fit(
            cropped,
            target_size,
            method=Image.Resampling.LANCZOS,
            centering=(0.5, 0.42),
        )
        image_stream = BytesIO()
        fitted.save(image_stream, format="PNG")
        image_stream.seek(0)

    _box(slide, x - 0.03, y - 0.03, width + 0.06, height + 0.06, fill=WHITE)
    slide.shapes.add_picture(
        image_stream,
        Inches(x),
        Inches(y),
        width=Inches(width),
        height=Inches(height),
    )


def _metric_summary() -> dict[str, Any]:
    dataset = generate_dataset()
    assert_quality(validate_dataset(dataset))
    kpis = build_daily_kpis(dataset.assets, dataset.generation, dataset.incidents)
    monthly: dict[str, dict[str, float]] = defaultdict(lambda: {"actual": 0.0, "forecast": 0.0})
    for row in kpis:
        key = row["generation_date"].strftime("%Y-%m")
        monthly[key]["actual"] += float(row["actual_generation_mwh"])
        monthly[key]["forecast"] += float(row["forecast_generation_mwh"])
    actual = sum(float(row["actual_generation_mwh"]) for row in kpis)
    forecast = sum(float(row["forecast_generation_mwh"]) for row in kpis)
    local_metrics = {
        "asset_rows": len(dataset.assets),
        "generation_rows": len(dataset.generation),
        "incident_rows": len(dataset.incidents),
        "kpi_rows": len(kpis),
        "total_generation_mwh": actual,
        "total_forecast_mwh": forecast,
        "variance_mwh": actual - forecast,
        "availability_pct": sum(float(row["availability_pct"]) for row in kpis) / len(kpis),
        "avoided_co2_tonnes": sum(float(row["avoided_co2_tonnes"]) for row in kpis),
        "monthly": dict(sorted(monthly.items())),
        "source": "Resultados locales deterministas validados",
    }
    if not REMOTE_EVIDENCE.exists():
        return local_metrics
    remote_metrics = json.loads(REMOTE_EVIDENCE.read_text(encoding="utf-8"))
    for key in ("asset_rows", "generation_rows", "incident_rows", "kpi_rows"):
        if remote_metrics[key] != local_metrics[key]:
            raise ValueError(
                f"Remote evidence mismatch for {key}: {remote_metrics[key]} != {local_metrics[key]}"
            )
    return remote_metrics


def _native_line_chart(slide: Slide, monthly: dict[str, dict[str, float]]) -> None:
    x, y, width, height = 0.95, 3.72, 7.4, 2.45
    axis = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x),
        Inches(y + height),
        Inches(x + width),
        Inches(y + height),
    )
    axis.line.color.rgb = LIGHT_LINE
    values = [series["actual"] for series in monthly.values()] + [
        series["forecast"] for series in monthly.values()
    ]
    minimum, maximum = min(values), max(values)
    months = list(monthly)
    for series_name, color in (("actual", TEAL), ("forecast", NAVY)):
        points: list[tuple[float, float]] = []
        for index, month in enumerate(months):
            px = x + index * width / (len(months) - 1)
            normalized = (monthly[month][series_name] - minimum) / (maximum - minimum)
            py = y + height - normalized * (height - 0.15)
            points.append((px, py))
        for start, end in pairwise(points):
            connector = slide.shapes.add_connector(
                MSO_CONNECTOR.STRAIGHT,
                Inches(start[0]),
                Inches(start[1]),
                Inches(end[0]),
                Inches(end[1]),
            )
            connector.line.color.rgb = color
            connector.line.width = Pt(2)
    _textbox(slide, "Real", 6.45, 3.45, 0.6, 0.2, size=10, color=TEAL, bold=True)
    _textbox(slide, "Prevista", 7.15, 3.45, 0.8, 0.2, size=10, color=NAVY, bold=True)
    _textbox(slide, months[0], x, y + height + 0.08, 0.9, 0.2, size=9, color=MID_GRAY)
    _textbox(
        slide,
        months[-1],
        x + width - 0.9,
        y + height + 0.08,
        0.9,
        0.2,
        size=9,
        color=MID_GRAY,
        align=PP_ALIGN.RIGHT,
    )


def _architecture_slide(presentation: PresentationType) -> None:
    slide = _base_slide(
        presentation,
        "Arquitectura",
        "Un flujo gobernado desde datos sintéticos hasta decisiones",
    )
    labels = [
        ("Datos\nsintéticos", RED),
        ("Tablas\nDelta", NAVY),
        ("Proceso\nserverless", NAVY),
        ("Metric View /\nsemantic view", TEAL),
        ("AI/BI\nDashboard", NAVY),
        ("Genie", ORANGE),
        ("Usuarios\nde negocio", NAVY),
    ]
    start_x = 0.55
    box_width = 1.55
    gap = 0.28
    for index, (label, color) in enumerate(labels):
        x = start_x + index * (box_width + gap)
        _box(slide, x, 2.65, box_width, 1.35, fill=GRAY, line=color)
        _textbox(
            slide,
            label,
            x + 0.08,
            2.82,
            box_width - 0.16,
            0.95,
            size=13,
            bold=True,
            align=PP_ALIGN.CENTER,
        )
        if index < len(labels) - 1:
            arrow = slide.shapes.add_shape(
                MSO_SHAPE.CHEVRON,
                Inches(x + box_width + 0.04),
                Inches(3.1),
                Inches(0.2),
                Inches(0.42),
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = RED
            arrow.line.fill.background()
    _textbox(
        slide,
        "Unity Catalog · definiciones reutilizables · infraestructura como código",
        2.15,
        4.7,
        9.0,
        0.5,
        size=17,
        color=TEAL,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    _notes(
        slide,
        "La misma definición gobernada alimenta dashboard y conversación.",
        "Pasemos de la arquitectura al modelo de datos.",
        "1:15",
    )


def build_legacy_presentation(metrics: dict[str, Any]) -> PresentationType:
    """Build the unused legacy deck retained for design reference."""
    presentation = Presentation()
    presentation.slide_width = SLIDE_WIDTH
    presentation.slide_height = SLIDE_HEIGHT

    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    red_block = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.22), SLIDE_HEIGHT
    )
    red_block.fill.solid()
    red_block.fill.fore_color.rgb = RED
    red_block.line.fill.background()
    _textbox(slide, "Renewable Operations\nIntelligence", 0.9, 1.35, 8.8, 1.45, size=32, bold=True)
    _textbox(
        slide,
        "De datos operativos a decisiones en lenguaje natural",
        0.93,
        3.0,
        8.5,
        0.55,
        size=19,
        color=MID_GRAY,
    )
    _textbox(
        slide, "Demo Databricks Free Edition", 0.93, 3.78, 5.8, 0.45, size=16, color=RED, bold=True
    )
    _box(slide, 10.05, 1.25, 2.25, 3.5, fill=GRAY, line=LIGHT_LINE)
    for index, (height, color) in enumerate(((1.35, NAVY), (2.05, TEAL), (2.7, RED))):
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(10.48 + index * 0.52),
            Inches(4.25 - height),
            Inches(0.3),
            Inches(height),
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = color
        bar.line.fill.background()
    _footer(slide, 1)
    _notes(
        slide,
        "Presentar la demo y recordar que todos los datos son sintéticos.",
        "Comencemos por el reto que resuelve.",
        "0:45",
    )

    slide = _base_slide(presentation, "Reto de negocio")
    challenges = [
        ("Información fragmentada", "Fuentes operativas aisladas y sin una lectura común."),
        ("KPIs inconsistentes", "La misma métrica puede calcularse de formas distintas."),
        ("Análisis manual", "Demasiado tiempo preparando consultas y reconciliando resultados."),
        (
            "Respuesta lenta",
            "Las desviaciones tardan en convertirse en una conversación accionable.",
        ),
    ]
    for index, (title, body) in enumerate(challenges):
        _card(slide, title, body, 0.75 + (index % 2) * 6.05, 1.5 + (index // 2) * 2.25, 5.75, 1.8)
    _notes(
        slide,
        "El problema no es solo acceder a datos, sino compartir definiciones.",
        "La propuesta combina gobierno, BI y conversación.",
        "1:00",
    )

    slide = _base_slide(presentation, "Propuesta de valor")
    values = [
        ("01", "Plataforma gobernada"),
        ("02", "Semantic layer reutilizable"),
        ("03", "Dashboard ejecutivo"),
        ("04", "Análisis conversacional con Genie"),
        ("05", "Despliegue como código"),
    ]
    for index, (number, label) in enumerate(values):
        y = 1.42 + index * 1.03
        _box(slide, 1.0, y, 11.0, 0.78, fill=GRAY)
        _textbox(slide, number, 1.22, y + 0.08, 0.65, 0.55, size=18, color=RED, bold=True)
        _textbox(slide, label, 2.28, y + 0.08, 8.65, 0.55, size=17, bold=True)
    _notes(
        slide,
        "Un único patrón sirve a operaciones, dirección y equipo de datos.",
        "Veamos cómo se conectan sus componentes.",
        "1:00",
    )

    _architecture_slide(presentation)

    slide = _base_slide(presentation, "Modelo de datos")
    model = [
        ("Assets", "10 instalaciones\ncapacidad · región · tecnología", 0.75, 1.6),
        ("Daily generation", "5.460 observaciones\nreal · prevista · disponibilidad", 6.95, 1.6),
        ("Incidents", "15 incidencias\nseveridad · downtime · estado", 0.75, 4.3),
        ("Daily KPIs", "variance · factor capacidad\ncoste/MWh · CO2 evitado", 6.95, 4.3),
    ]
    for title, body, x, y in model:
        _card(slide, title, body, x, y, 5.6, 1.6, accent=TEAL if title == "Daily KPIs" else RED)
    for start, end in (
        ((6.35, 2.4), (6.9, 2.4)),
        ((9.75, 3.2), (9.75, 4.25)),
        ((6.35, 5.1), (6.9, 5.1)),
    ):
        connector = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            Inches(start[0]),
            Inches(start[1]),
            Inches(end[0]),
            Inches(end[1]),
        )
        connector.line.color.rgb = MID_GRAY
        connector.line.width = Pt(1.5)
    _notes(
        slide,
        "El KPI diario concentra la semántica necesaria para consumo seguro.",
        "Sobre ese modelo se definen métricas gobernadas.",
        "1:00",
    )

    slide = _base_slide(presentation, "Métricas gobernadas")
    metric_names = [
        "Generación",
        "Previsión",
        "Desviación",
        "Disponibilidad",
        "Factor de capacidad",
        "Coste por MWh",
        "CO2 evitado",
    ]
    for index, name in enumerate(metric_names):
        row, column = divmod(index, 4)
        x, y = 0.7 + column * 3.15, 1.6 + row * 2.15
        _box(slide, x, y, 2.85, 1.55, fill=GRAY)
        _textbox(
            slide, name, x + 0.15, y + 0.22, 2.55, 0.4, size=15, bold=True, align=PP_ALIGN.CENTER
        )
        rule = "Definición única · unidad visible"
        if name == "Desviación":
            rule = "Real - prevista · negativa = desfavorable"
        elif name == "Coste por MWh":
            rule = "EUR / generación real · cero → NULL"
        _textbox(
            slide,
            rule,
            x + 0.18,
            y + 0.72,
            2.5,
            0.48,
            size=10,
            color=MID_GRAY,
            align=PP_ALIGN.CENTER,
        )
    _textbox(
        slide,
        "Calendario natural · periodo siempre explícito",
        3.55,
        6.05,
        6.2,
        0.45,
        size=16,
        color=TEAL,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    _notes(
        slide,
        "Las reglas evitan que dashboard y Genie respondan con lógicas distintas.",
        "Apliquemos esas métricas a una vista ejecutiva.",
        "1:00",
    )

    slide = _base_slide(
        presentation,
        "Dashboard ejecutivo",
        "Dashboard AI/BI publicado · paleta validada en tema oscuro",
    )
    _real_screenshot(slide, DASHBOARD_SCREENSHOT, 0.68, 1.25, 11.98, 5.35)
    _textbox(
        slide,
        "Pantalla real · ambas series y barras visibles · filtros globales activos",
        0.82,
        6.65,
        11.65,
        0.24,
        size=10,
        color=MID_GRAY,
        align=PP_ALIGN.CENTER,
    )
    _notes(
        slide,
        "Mostrar el dashboard real corregido y explicar las dos series y los filtros.",
        "La misma capa sirve a una conversación en lenguaje natural.",
        "1:45",
        "Abrir el dashboard real, filtrar el mes adverso y bajar a región.",
    )

    slide = _base_slide(
        presentation,
        "Conversación con Genie",
        "Agent desplegado · fuente semántica gobernada",
    )
    _real_screenshot(slide, GENIE_SCREENSHOT, 0.68, 1.25, 11.98, 5.35)
    _textbox(
        slide,
        "Fuente autorizada: workspace.renewable_operations_demo.gg_renewable_operations_metrics",
        0.82,
        6.65,
        11.65,
        0.24,
        size=10,
        color=TEAL,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    _notes(
        slide,
        "Genie acelera la exploración sin abandonar las definiciones acordadas.",
        "Convirtamos la conversación en un guion repetible.",
        "1:15",
        "Ejecutar pregunta inicial y follow-up de instalación e incidencias.",
    )

    slide = _base_slide(presentation, "Guion de demo")
    steps = [
        "Detectar\ndesviación",
        "Analizar\nregión",
        "Bajar a\ninstalación",
        "Relacionar\nincidencias",
        "Preguntar en\nlenguaje natural",
        "Convertir en\nacción",
    ]
    for index, step in enumerate(steps):
        x = 0.55 + index * 2.08
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(x + 0.55), Inches(2.3), Inches(0.72), Inches(0.72)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = RED if index in {0, 5} else NAVY
        circle.line.fill.background()
        _textbox(
            slide,
            str(index + 1),
            x + 0.67,
            2.44,
            0.48,
            0.38,
            size=15,
            color=WHITE,
            bold=True,
            align=PP_ALIGN.CENTER,
        )
        _textbox(slide, step, x, 3.28, 1.85, 0.85, size=13, bold=True, align=PP_ALIGN.CENTER)
        if index < len(steps) - 1:
            connector = slide.shapes.add_connector(
                MSO_CONNECTOR.STRAIGHT,
                Inches(x + 1.27),
                Inches(2.66),
                Inches(x + 2.63),
                Inches(2.66),
            )
            connector.line.color.rgb = LIGHT_LINE
            connector.line.width = Pt(2)
    _textbox(
        slide,
        "Una historia de 12-15 minutos: del síntoma a una decisión trazable",
        2.25,
        5.25,
        8.8,
        0.55,
        size=17,
        color=TEAL,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    _notes(
        slide,
        "El guion evita una sucesión de pantallas y mantiene una narrativa de decisión.",
        "Ahora resumamos el impacto esperado.",
        "0:50",
    )

    slide = _base_slide(presentation, "Valor para negocio")
    benefits = [
        "Menor tiempo hasta el insight",
        "Autoservicio guiado",
        "Coherencia métrica",
        "Trazabilidad",
        "Menor carga de consultas repetitivas",
        "Patrón escalable",
    ]
    for index, benefit in enumerate(benefits):
        row, column = divmod(index, 3)
        _card(
            slide,
            benefit,
            "Resultado esperable; sujeto a validación en el contexto corporativo.",
            0.7 + column * 4.15,
            1.5 + row * 2.45,
            3.85,
            1.85,
            accent=TEAL if index in {1, 2, 3} else RED,
        )
    _notes(
        slide,
        "Hablar de resultados esperables, no de ahorros garantizados.",
        "La confianza depende de controles técnicos visibles.",
        "1:00",
    )

    slide = _base_slide(presentation, "Controles y calidad")
    controls = [
        ("Infraestructura como código", "Bundle validable y reproducible"),
        ("Tests", "29 tests · 92,0 % cobertura"),
        ("Data quality", "Claves, rangos, conteos y coherencia"),
        ("Idempotencia", "Segunda ejecución sin duplicados"),
        ("Unity Catalog", "Namespace aislado y comentado"),
        ("Aislamiento Genie", "Solo semantic layer autorizada"),
    ]
    for index, (title, body) in enumerate(controls):
        row, column = divmod(index, 2)
        _card(slide, title, body, 0.75 + column * 6.1, 1.4 + row * 1.7, 5.8, 1.35, accent=TEAL)
    contract_source = "desplegado" if "remoto" in metrics["source"] else "local"
    _textbox(
        slide,
        (
            f"Contrato {contract_source}: {metrics['asset_rows']} assets · "
            f"{metrics['generation_rows']:,} días-activo · "
            f"{metrics['incident_rows']} incidencias"
        ),
        1.25,
        6.35,
        10.8,
        0.4,
        size=14,
        color=NAVY,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    _notes(
        slide,
        "La demo incluye evidencia técnica, no solo visualizaciones.",
        "Cerramos con el camino de industrialización.",
        "1:00",
        "Mostrar el run y el informe de calidad.",
    )

    slide = _base_slide(presentation, "Próximos pasos")
    next_steps = [
        "Incorporar datos corporativos gobernados",
        "Separar entornos y permisos",
        "Añadir dominios y acuerdos de calidad",
        "Validar seguridad de extremo a extremo",
        "Medir adopción y calidad de respuestas",
        "Industrializar CI/CD",
    ]
    for index, step in enumerate(next_steps):
        y = 1.35 + index * 0.82
        _textbox(slide, f"{index + 1:02d}", 1.0, y, 0.55, 0.45, size=15, color=RED, bold=True)
        _textbox(slide, step, 1.75, y, 8.6, 0.45, size=16, bold=True)
    _box(slide, 10.4, 1.55, 2.15, 3.6, fill=GRAY)
    _textbox(slide, "Estado", 10.75, 1.9, 1.45, 0.4, size=16, bold=True, align=PP_ALIGN.CENTER)
    remote_validated = "remoto" in metrics["source"]
    _textbox(
        slide,
        "Despliegue remoto\nvalidado" if remote_validated else "Validación local\ncompletada",
        10.7,
        2.65,
        1.55,
        0.85,
        size=15,
        color=TEAL,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    _textbox(
        slide,
        (
            "Dos ejecuciones\nidempotentes"
            if remote_validated
            else "Despliegue remoto\nse refleja en el\ninforme final"
        ),
        10.72,
        3.75,
        1.5,
        1.0,
        size=11,
        color=MID_GRAY,
        align=PP_ALIGN.CENTER,
    )
    _notes(
        slide,
        "La demo prueba el patrón; industrializar requiere gobierno y medición continua.",
        "Abrir preguntas.",
        "0:55",
    )
    return presentation


def build_presentation(metrics: dict[str, Any]) -> PresentationType:
    """Build the client-ready business presentation."""
    presentation = Presentation()
    presentation.slide_width = SLIDE_WIDTH
    presentation.slide_height = SLIDE_HEIGHT
    presentation.core_properties.title = (
        "Inteligencia operativa para energía renovable con Databricks Genie"
    )
    presentation.core_properties.subject = (
        "Prueba de viabilidad de AI/BI Dashboard, Genie Agent y Genie One"
    )
    presentation.core_properties.author = "GreenGrid Energy demo"

    total_generation = f"{metrics['total_generation_mwh'] / 1_000_000:.1f} M".replace(".", ",")
    variance = f"{metrics['variance_mwh'] / 1_000:.1f} mil".replace(".", ",")
    availability = f"{metrics['availability_pct']:.2f} %".replace(".", ",")
    co2 = f"{metrics['avoided_co2_tonnes'] / 1_000_000:.2f} M".replace(".", ",")

    # 1 · Cover
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = NAVY
    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.16), SLIDE_HEIGHT
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = RED
    accent.line.fill.background()
    _pill(
        slide,
        "VIABILIDAD E IMPLEMENTACIÓN · DATABRICKS GENIE",
        0.78,
        0.62,
        4.25,
        fill=RED,
    )
    _textbox(
        slide,
        "Inteligencia operativa\npara energía renovable",
        0.78,
        1.35,
        5.35,
        1.55,
        size=31,
        color=WHITE,
        bold=True,
        valign=MSO_ANCHOR.TOP,
    )
    _textbox(
        slide,
        "Genie One aplicado a generación, fiabilidad y decisiones operativas",
        0.82,
        3.05,
        5.1,
        0.72,
        size=18,
        color=LIGHT_LINE,
        valign=MSO_ANCHOR.TOP,
    )
    _textbox(
        slide,
        "AI/BI Dashboard · Genie Agent · Genie One · Automation Bundle",
        0.82,
        4.05,
        5.2,
        0.55,
        size=12,
        color=TEAL,
        bold=True,
    )
    _box(
        slide,
        6.28,
        0.72,
        6.35,
        5.9,
        fill=RGBColor(0x23, 0x30, 0x40),
        line=MID_GRAY,
    )
    _textbox(
        slide,
        "ARQUITECTURA DEL DEMO",
        6.65,
        1.02,
        5.58,
        0.32,
        size=12,
        color=LIGHT_LINE,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    _box(slide, 7.02, 1.52, 4.85, 0.7, fill=RED, line=RED)
    _textbox(
        slide,
        "GENIE ONE",
        7.2,
        1.59,
        4.49,
        0.25,
        size=13,
        color=WHITE,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    _textbox(
        slide,
        "Entrada única para negocio",
        7.2,
        1.91,
        4.49,
        0.25,
        size=10,
        color=WHITE,
        align=PP_ALIGN.CENTER,
    )
    _arrow(slide, (9.45, 2.23), (9.45, 2.48), color=RED)
    experiences = [
        (
            "AI/BI DASHBOARD",
            "KPIs · filtros · visuales",
            6.62,
            TEAL,
        ),
        (
            "GENIE AGENT",
            "Preguntas · SQL · benchmark",
            9.54,
            ORANGE,
        ),
    ]
    for title, body, x, color in experiences:
        _box(slide, x, 2.5, 2.7, 0.92, fill=NAVY, line=color)
        _textbox(
            slide,
            title,
            x + 0.16,
            2.62,
            2.38,
            0.27,
            size=11,
            color=color,
            bold=True,
            align=PP_ALIGN.CENTER,
        )
        _textbox(
            slide,
            body,
            x + 0.16,
            3.02,
            2.38,
            0.22,
            size=8.5,
            color=LIGHT_LINE,
            align=PP_ALIGN.CENTER,
        )
    _arrow(slide, (9.45, 3.44), (9.45, 3.7), color=TEAL)
    cover_layers = [
        (
            "UNITY CATALOG + SEMÁNTICA",
            "KPI diario · Semantic View · Metric View",
            3.72,
            TEAL,
        ),
        (
            "SQL WAREHOUSE + WORKFLOW",
            "Consulta y procesamiento serverless",
            4.73,
            ORANGE,
        ),
        (
            "DATOS SINTÉTICOS DE OPERACIÓN ENERGÉTICA",
            "Generación · previsión · disponibilidad · incidencias",
            5.72,
            LIGHT_LINE,
        ),
    ]
    for title, body, y, color in cover_layers:
        _box(slide, 7.02, y, 4.85, 0.72, fill=NAVY, line=color)
        _textbox(
            slide,
            title,
            7.2,
            y + 0.1,
            4.45,
            0.24,
            size=10.5,
            color=color,
            bold=True,
            align=PP_ALIGN.CENTER,
        )
        _textbox(
            slide,
            body,
            7.2,
            y + 0.4,
            4.45,
            0.2,
            size=8.3,
            color=LIGHT_LINE,
            align=PP_ALIGN.CENTER,
        )
    _arrow(slide, (9.45, 4.45), (9.45, 4.7), color=TEAL)
    _arrow(slide, (9.45, 5.46), (9.45, 5.69), color=ORANGE)
    _textbox(slide, FOOTER, 0.82, 7.1, 7.4, 0.2, size=9, color=LIGHT_LINE)
    _textbox(slide, "1", 12.08, 7.1, 0.55, 0.2, size=9, color=LIGHT_LINE)
    _notes(
        slide,
        "Abrir con el resultado: una experiencia de negocio completa sobre una base gobernada.",
        "Primero resumimos qué demuestra realmente el entregable.",
        "0:45",
    )

    # 2 · Executive promise
    slide = _base_slide(
        presentation,
        "Qué demuestra esta demo",
        "Una historia completa: observar, preguntar, validar y actuar",
    )
    story = [
        ("01 · OBSERVAR", "KPIs ejecutivos y fiabilidad operativa", RED),
        ("02 · PREGUNTAR", "Análisis conversacional sobre la misma semántica", ORANGE),
        ("03 · CONFIAR", "SQL de referencia, benchmarks y trazabilidad", TEAL),
        ("04 · ESCALAR", "Entrega reproducible y extensible con Apps", NAVY),
    ]
    for index, (title, body, color) in enumerate(story):
        x = 0.72 + index * 3.08
        _box(slide, x, 1.58, 2.78, 3.8, fill=GRAY, line=color)
        _textbox(slide, title, x + 0.18, 1.88, 2.42, 0.5, size=13, color=color, bold=True)
        _textbox(
            slide,
            body,
            x + 0.18,
            2.58,
            2.4,
            1.28,
            size=17,
            bold=True,
            valign=MSO_ANCHOR.TOP,
        )
        _textbox(
            slide,
            (
                "Dashboard publicado"
                if index == 0
                else "Genie Agent real"
                if index == 1
                else "5/5 benchmarks"
                if index == 2
                else "GitHub + Bundle"
            ),
            x + 0.18,
            4.45,
            2.4,
            0.38,
            size=10,
            color=MID_GRAY,
        )
    _textbox(
        slide,
        "El activo principal no es una pantalla: es un patrón reutilizable de decisión gobernada.",
        1.2,
        5.9,
        10.9,
        0.52,
        size=18,
        color=TEAL,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    _notes(
        slide,
        "Aclarar que cada promesa está respaldada por un recurso desplegado o una prueba real.",
        "Conectamos la propuesta con un problema de negocio reconocible.",
        "1:00",
    )

    # 3 · Business challenge
    slide = _dark_slide(
        presentation,
        "El reto: del dato operativo a una decisión compartida",
        "La dificultad no es solo consultar datos; es acordar significado y acelerar la respuesta",
    )
    challenges = [
        (
            "Señales dispersas",
            "Generación, previsión, disponibilidad e incidencias sin una lectura común.",
        ),
        (
            "Métricas discutibles",
            "Distintos cálculos para desviación, coste o factor de capacidad.",
        ),
        (
            "Exploración lenta",
            "Cada nueva pregunta vuelve a depender de un analista o una consulta ad hoc.",
        ),
        ("Acción desconectada", "El insight no siempre llega al flujo operativo adecuado."),
    ]
    for index, (title, body) in enumerate(challenges):
        row, column = divmod(index, 2)
        x, y = 0.82 + column * 6.12, 1.58 + row * 2.34
        _box(slide, x, y, 5.72, 1.92, fill=RGBColor(0x23, 0x30, 0x40), line=MID_GRAY)
        _textbox(
            slide, f"0{index + 1}", x + 0.22, y + 0.2, 0.6, 0.45, size=16, color=RED, bold=True
        )
        _textbox(slide, title, x + 1.02, y + 0.2, 4.3, 0.42, size=16, color=WHITE, bold=True)
        _textbox(
            slide,
            body,
            x + 1.02,
            y + 0.75,
            4.25,
            0.8,
            size=12,
            color=LIGHT_LINE,
            valign=MSO_ANCHOR.TOP,
        )
    _notes(
        slide,
        "Enmarcar la demo como solución al ciclo completo de decisión y no como chatbot aislado.",
        "Mostramos el caso de uso que hemos implementado.",
        "1:00",
    )

    # 4 · Use case and real KPIs
    slide = _base_slide(
        presentation,
        "Caso de uso: operación de activos de energía renovable",
        "18 meses · 10 instalaciones · datos exclusivamente sintéticos",
    )
    _kpi_card(slide, total_generation, "Generación real (MWh)", 0.72, 1.42, 2.75, accent=TEAL)
    _kpi_card(slide, variance, "Desviación vs. previsión (MWh)", 3.61, 1.42, 2.75, accent=RED)
    _kpi_card(slide, availability, "Disponibilidad media", 6.5, 1.42, 2.75, accent=TEAL)
    _kpi_card(slide, co2, "CO₂ evitado (toneladas)", 9.39, 1.42, 2.75, accent=ORANGE)
    _textbox(
        slide,
        (
            "GreenGrid Energy opera instalaciones solares, eólicas e hidráulicas. "
            "La solución identifica desviaciones de generación y fiabilidad para "
            "priorizar dónde investigar y actuar."
        ),
        0.92,
        2.73,
        11.5,
        0.42,
        size=11.5,
        color=MID_GRAY,
        align=PP_ALIGN.CENTER,
    )
    questions = [
        ("Dirección", "¿Dónde está la desviación y cuál es su impacto?"),
        ("Operaciones", "¿Qué instalación combina baja disponibilidad e incidencias?"),
        ("Control", "¿Qué tecnología presenta mayor coste por MWh?"),
    ]
    for index, (role, question) in enumerate(questions):
        x = 0.82 + index * 4.12
        _card(slide, role, question, x, 3.2, 3.72, 1.75, accent=(RED, TEAL, ORANGE)[index])
    _textbox(
        slide,
        "Una misma capa semántica responde tanto al dashboard como a Genie.",
        2.0,
        5.58,
        9.35,
        0.48,
        size=18,
        color=NAVY,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    _notes(
        slide,
        "Usar los KPIs reales como ancla y presentar las preguntas por perfil de negocio.",
        "Situamos ahora cada experiencia dentro de la familia Genie.",
        "1:15",
    )

    # 5 · Genie family
    slide = _dark_slide(
        presentation,
        "La familia Genie: tres experiencias, una base gobernada",
        "Terminología vigente: Genie Agents era anteriormente Genie Spaces",
    )
    genie_family = [
        (
            "Genie One",
            "Interfaz simplificada para negocio",
            "Chat unificado · dashboards · Apps · dominios",
            RED,
        ),
        (
            "Genie Agents",
            "Especialistas por dominio",
            "Datos curados · instrucciones · SQL de referencia · benchmarks",
            TEAL,
        ),
        (
            "Genie Code",
            "Asistente para equipos técnicos",
            "Código · pipelines · dashboards · tareas agentic",
            ORANGE,
        ),
    ]
    for index, (name, audience, body, color) in enumerate(genie_family):
        x = 0.82 + index * 4.12
        _box(slide, x, 1.65, 3.72, 3.98, fill=RGBColor(0x23, 0x30, 0x40), line=color)
        _pill(slide, name, x + 0.25, 1.95, 1.65, fill=color)
        _textbox(slide, audience, x + 0.25, 2.65, 3.1, 0.72, size=18, color=WHITE, bold=True)
        _textbox(
            slide,
            body,
            x + 0.25,
            3.62,
            3.1,
            1.1,
            size=13,
            color=LIGHT_LINE,
            valign=MSO_ANCHOR.TOP,
        )
    _textbox(
        slide,
        "Este demo despliega un Genie Agent y lo hace descubrible desde Genie One.",
        1.4,
        6.1,
        10.55,
        0.42,
        size=17,
        color=TEAL,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    _source_note(slide, "Fuente: docs.databricks.com/aws/en/genie/ · consultado 24/07/2026")
    _notes(
        slide,
        (
            "Evitar mezclar Genie One con el Agent: One es el punto de consumo; "
            "el Agent contiene contexto."
        ),
        "Veamos la experiencia real que recibe un usuario de negocio.",
        "1:20",
    )

    # 6 · Genie One real
    slide = _base_slide(
        presentation,
        "Genie One: un único punto de entrada para negocio",
        "Sin notebooks, clusters ni SQL como punto de partida",
    )
    _real_screenshot(slide, GENIE_ONE_HOME_SCREENSHOT, 0.72, 1.3, 8.25, 5.22)
    benefits = [
        ("Preguntar", "Lenguaje natural sobre datos gobernados."),
        ("Descubrir", "Agents, dashboards y dominios compartidos."),
        ("Consumir", "Apps y tareas sin cambiar de experiencia."),
        ("Movilidad", "Experiencia disponible también en móvil."),
    ]
    for index, (title, body) in enumerate(benefits):
        y = 1.35 + index * 1.27
        _card(slide, title, body, 9.28, y, 3.35, 1.03, accent=(RED, TEAL, ORANGE, NAVY)[index])
    _source_note(slide, "Fuente: docs.databricks.com/aws/en/genie-one/")
    _notes(
        slide,
        "Señalar que esta pantalla reduce la exposición a conceptos técnicos del workspace.",
        "Bajamos de la experiencia a la jerarquía de componentes.",
        "1:20",
        "Abrir /one y mostrar Agents y Dashboards.",
    )

    # 7 · Hierarchy
    slide = _base_slide(
        presentation,
        "Jerarquía del demo en Databricks",
        "Qué contiene cada nivel y quién lo consume",
    )
    levels = [
        ("CUENTA", "Identidad · usuarios · acceso", 0.8, 1.4, 11.75, 0.62, NAVY),
        (
            "WORKSPACE",
            "Free Edition · compute serverless · SQL warehouse",
            1.18,
            2.18,
            11.0,
            0.7,
            RED,
        ),
        (
            "UNITY CATALOG",
            "catalog workspace  >  schema renewable_operations_demo",
            1.58,
            3.06,
            10.2,
            0.76,
            TEAL,
        ),
    ]
    for label, body, x, y, width, height, color in levels:
        _box(slide, x, y, width, height, fill=GRAY, line=color)
        _textbox(
            slide, label, x + 0.18, y + 0.1, 1.6, height - 0.2, size=11, color=color, bold=True
        )
        _textbox(slide, body, x + 2.0, y + 0.1, width - 2.2, height - 0.2, size=13, bold=True)
    components = [
        ("Delta tables", "assets · generación · incidencias"),
        ("Capa semántica", "daily KPI · semantic view · Metric View"),
        ("Servicio", "SQL Warehouse · workflow serverless"),
        ("Experiencias", "AI/BI Dashboard · Genie Agent"),
    ]
    for index, (title, body) in enumerate(components):
        x = 0.72 + index * 3.08
        _box(slide, x, 4.18, 2.78, 1.35, fill=WHITE, line=(NAVY, TEAL, ORANGE, RED)[index])
        _textbox(slide, title, x + 0.15, 4.35, 2.48, 0.35, size=13, bold=True)
        _textbox(
            slide,
            body,
            x + 0.15,
            4.78,
            2.48,
            0.45,
            size=10,
            color=MID_GRAY,
            valign=MSO_ANCHOR.TOP,
        )
    _arrow(slide, (6.67, 3.83), (6.67, 4.12), color=TEAL)
    _box(slide, 2.45, 5.9, 8.45, 0.62, fill=NAVY, line=NAVY)
    _textbox(
        slide,
        "GENIE ONE · consumo por usuarios de negocio  |  APP operativa · extensión opcional",
        2.65,
        6.02,
        8.05,
        0.35,
        size=12,
        color=WHITE,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    _notes(
        slide,
        (
            "Explicar la jerarquía desde identidad hasta consumo y distinguir "
            "objetos de datos de experiencias."
        ),
        "La siguiente lámina lo convierte en una visión visual de extremo a extremo.",
        "1:30",
    )

    # 8 · Flat architecture flow
    slide = _base_slide(
        presentation,
        "De los datos operativos a la acción",
        "Componentes del demo y flujo de valor · diagrama plano y editable",
    )
    stages = [
        (
            "01",
            "DATOS",
            "Generación real y prevista\n\nDisponibilidad\n\nIncidencias y costes",
            NAVY,
        ),
        (
            "02",
            "PREPARAR",
            "Tablas Delta\n\nKPI diario\n\nWorkflow serverless",
            ORANGE,
        ),
        (
            "03",
            "GOBERNAR",
            "Unity Catalog\n\nSemantic View\n\nMetric View\n\nUnidades y periodo",
            TEAL,
        ),
        (
            "04",
            "ANALIZAR",
            "AI/BI Dashboard\n\nGenie Agent\n\nSQL de referencia\n\nBenchmarks",
            RED,
        ),
        (
            "05",
            "DECIDIR",
            "Genie One\n\nObservar\n\nPreguntar\n\nPriorizar",
            NAVY,
        ),
    ]
    for index, (number, title, body, color) in enumerate(stages):
        x = 0.58 + index * 2.55
        _box(slide, x, 1.58, 2.22, 4.7, fill=GRAY, line=color)
        marker = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(x),
            Inches(1.58),
            Inches(2.22),
            Inches(0.09),
        )
        marker.fill.solid()
        marker.fill.fore_color.rgb = color
        marker.line.fill.background()
        _textbox(
            slide,
            number,
            x + 0.18,
            1.88,
            0.52,
            0.4,
            size=16,
            color=color,
            bold=True,
        )
        _textbox(
            slide,
            title,
            x + 0.18,
            2.38,
            1.84,
            0.4,
            size=15,
            color=NAVY,
            bold=True,
        )
        _textbox(
            slide,
            body,
            x + 0.18,
            3.05,
            1.84,
            2.55,
            size=11,
            color=MID_GRAY,
            valign=MSO_ANCHOR.TOP,
        )
        if index < len(stages) - 1:
            _arrow(slide, (x + 2.24, 3.92), (x + 2.49, 3.92), color=color, width=2.4)
    _box(slide, 1.2, 6.52, 10.9, 0.42, fill=NAVY, line=NAVY)
    _textbox(
        slide,
        "DESPLEGADO HOY: DATOS → PREPARACIÓN → GOBIERNO → ANÁLISIS → DECISIÓN",
        1.42,
        6.58,
        10.46,
        0.25,
        size=10.5,
        color=WHITE,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    _notes(
        slide,
        (
            "Recorrer el diagrama plano de izquierda a derecha y explicar "
            "la función de cada componente."
        ),
        "Aterrizamos esa visión en el modelo de datos gobernado.",
        "1:00",
    )

    # 9 · Data foundation
    slide = _base_slide(
        presentation,
        "Fundación gobernada y reutilizable",
        "La calidad de Genie empieza antes del prompt",
    )
    foundation = [
        ("01", "Datos Delta", "10 assets · 5.460 días-activo · 15 incidencias", NAVY),
        ("02", "KPI diario", "Real · prevista · disponibilidad · coste · CO₂", RED),
        ("03", "Semantic View", "Fallback tabular y dataset del dashboard", ORANGE),
        ("04", "Metric View", "Medidas gobernadas para Genie y reutilización", TEAL),
    ]
    for index, (number, title, body, color) in enumerate(foundation):
        y = 1.38 + index * 1.23
        _box(slide, 0.82, y, 7.2, 0.98, fill=GRAY, line=color)
        _textbox(slide, number, 1.02, y + 0.17, 0.58, 0.45, size=16, color=color, bold=True)
        _textbox(slide, title, 1.75, y + 0.12, 1.75, 0.35, size=14, bold=True)
        _textbox(slide, body, 3.55, y + 0.12, 4.1, 0.5, size=11, color=MID_GRAY)
        if index < len(foundation) - 1:
            _arrow(slide, (4.35, y + 0.99), (4.35, y + 1.18), color=color)
    _box(slide, 8.48, 1.42, 4.1, 4.82, fill=NAVY, line=NAVY)
    _textbox(slide, "Contrato semántico", 8.85, 1.8, 3.35, 0.52, size=19, color=WHITE, bold=True)
    semantics = [
        "Desviación = real - prevista",
        "Importes en EUR",
        "Generación en MWh",
        "Periodo siempre explícito",
        "Ceros sensibles → NULL",
        "Instalación ≠ equipo operador",
    ]
    for index, item in enumerate(semantics):
        _textbox(slide, "✓", 8.9, 2.55 + index * 0.52, 0.3, 0.3, size=13, color=TEAL, bold=True)
        _textbox(slide, item, 9.35, 2.51 + index * 0.52, 2.72, 0.34, size=11, color=LIGHT_LINE)
    _notes(
        slide,
        (
            "Subrayar que la exactitud no depende solo del modelo: depende de "
            "métricas y contexto curados."
        ),
        "Mostramos cómo se consume esa base en el dashboard.",
        "1:10",
    )

    # 10 · Executive dashboard
    slide = _base_slide(
        presentation,
        "AI/BI Dashboard · visión ejecutiva",
        "Filtros globales y KPIs gobernados",
    )
    _real_screenshot(slide, DASHBOARD_SCREENSHOT, 0.68, 1.25, 11.98, 5.35)
    _box(slide, 0.78, 6.14, 11.78, 0.55, fill=WHITE, line=LIGHT_LINE)
    _pill(slide, "SEÑAL", 0.86, 6.32, 0.85, fill=RED)
    _textbox(
        slide,
        "La generación real queda por debajo de la previsión en el periodo analizado.",
        1.88,
        6.31,
        10.3,
        0.32,
        size=11,
        color=NAVY,
        bold=True,
    )
    _notes(
        slide,
        "Mostrar KPIs, tendencia, tecnología, región y filtros; evitar leer cada gráfico.",
        "Bajamos a la fiabilidad operativa.",
        "1:30",
        "Filtrar por tecnología o región en el dashboard real.",
    )

    # 11 · Operational dashboard
    slide = _base_slide(
        presentation,
        "AI/BI Dashboard · fiabilidad operativa",
        "Nombres de instalación explícitos y dimensiones separadas",
    )
    _real_screenshot(slide, DASHBOARD_RELIABILITY_SCREENSHOT, 0.68, 1.25, 11.98, 5.35)
    _box(slide, 0.78, 6.14, 11.78, 0.55, fill=WHITE, line=LIGHT_LINE)
    _pill(slide, "FOCO", 0.86, 6.32, 0.85, fill=TEAL)
    _textbox(
        slide,
        "El parque eólico concentra el mayor downtime agregado y requiere investigación.",
        1.88,
        6.31,
        10.3,
        0.32,
        size=11,
        color=NAVY,
        bold=True,
    )
    _notes(
        slide,
        "Enseñar la transición de la visión ejecutiva a instalaciones, incidencias y coste.",
        "La siguiente pregunta profundiza sin abandonar la semántica.",
        "1:20",
        "Abrir la pestaña Operational Reliability.",
    )

    # 12 · Genie conversation
    slide = _base_slide(
        presentation,
        "Genie Agent · de una pregunta a una explicación",
        "Conversación real · SQL, tabla, gráfico y narrativa en una misma respuesta",
    )
    _real_screenshot(slide, GENIE_CONVERSATION_SCREENSHOT, 0.68, 1.25, 11.98, 5.35)
    _box(slide, 0.78, 6.14, 11.78, 0.55, fill=WHITE, line=LIGHT_LINE)
    _textbox(
        slide,
        "Hallazgo real: Nerea Insular · 95,27 % disponibilidad · 2 incidencias · 274 h de parada",
        0.88,
        6.34,
        11.55,
        0.3,
        size=11,
        color=TEAL,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    _notes(
        slide,
        (
            "Leer solo el hallazgo y señalar que el usuario puede inspeccionar "
            "SQL, tabla y visualización."
        ),
        "Descomponemos cómo se genera una respuesta confiable.",
        "1:40",
        "Abrir la conversación capturada y mostrar Show code.",
    )

    # 13 · Answer anatomy
    slide = _dark_slide(
        presentation,
        "Anatomía de una respuesta confiable",
        "Qué ocurre entre la pregunta y el insight",
    )
    stages = [
        ("1", "Pregunta", "Lenguaje y objetivo de negocio", RED),
        ("2", "Contexto", "Instrucciones + SQL de referencia", ORANGE),
        ("3", "Semántica", "Metric View + definiciones", TEAL),
        ("4", "Ejecución", "SQL Warehouse serverless", NAVY),
        ("5", "Respuesta", "Narrativa + tabla + gráfico", RED),
    ]
    for index, (number, title, body, color) in enumerate(stages):
        x = 0.48 + index * 2.55
        _box(slide, x, 2.0, 2.15, 2.38, fill=RGBColor(0x23, 0x30, 0x40), line=color)
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(x + 0.73), Inches(1.62), Inches(0.68), Inches(0.68)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = color
        circle.line.fill.background()
        _textbox(
            slide,
            number,
            x + 0.85,
            1.75,
            0.44,
            0.34,
            size=14,
            color=WHITE,
            bold=True,
            align=PP_ALIGN.CENTER,
        )
        _textbox(slide, title, x + 0.18, 2.5, 1.78, 0.4, size=15, color=WHITE, bold=True)
        _textbox(
            slide,
            body,
            x + 0.18,
            3.08,
            1.78,
            0.78,
            size=11,
            color=LIGHT_LINE,
            valign=MSO_ANCHOR.TOP,
        )
        if index < len(stages) - 1:
            _arrow(slide, (x + 2.16, 3.16), (x + 2.52, 3.16), color=RED)
    _box(slide, 1.25, 5.25, 10.85, 0.82, fill=RGBColor(0x23, 0x30, 0x40), line=TEAL)
    _textbox(
        slide,
        (
            "La respuesta queda trazada hasta workspace.renewable_operations_demo."
            "gg_renewable_operations_metrics"
        ),
        1.55,
        5.45,
        10.25,
        0.38,
        size=13,
        color=TEAL,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    _notes(
        slide,
        (
            "Explicar que Genie no responde desde conocimiento genérico: "
            "genera SQL sobre una fuente autorizada."
        ),
        "La confianza se verifica de forma repetible con benchmarks.",
        "1:20",
    )

    # 14 · Benchmark evidence
    slide = _base_slide(
        presentation,
        "Calidad medible · benchmarks del Genie Agent",
        "Evaluación real ejecutada el 24/07/2026",
    )
    _real_screenshot(slide, GENIE_BENCHMARK_SCREENSHOT, 0.68, 1.25, 11.98, 5.35)
    _box(slide, 0.78, 6.14, 11.78, 0.55, fill=WHITE, line=LIGHT_LINE)
    _pill(slide, "RESULTADO REAL", 0.86, 6.31, 1.52, fill=TEAL)
    _textbox(
        slide,
        "100 % accurate · 5/5 preguntas · comparación con SQL esperado",
        2.6,
        6.3,
        9.45,
        0.32,
        size=12,
        color=NAVY,
        bold=True,
    )
    _notes(
        slide,
        "Mostrar que la calidad puede medirse y revisarse por pregunta, SQL y ground truth.",
        "Resumimos el ciclo de mejora continua.",
        "1:30",
        "Abrir Benchmark y seleccionar una evaluación.",
    )

    # 15 · Quality loop
    slide = _base_slide(
        presentation,
        "El ciclo operativo de calidad",
        "Configurar · probar · observar · corregir · volver a medir",
    )
    quality = [
        ("CURAR", "Fuentes, columnas, métricas y terminología"),
        ("ENSEÑAR", "Instrucciones y SQL de referencia"),
        ("EVALUAR", "Benchmarks frecuentes y variantes de preguntas"),
        ("MONITORIZAR", "Feedback, revisiones y consultas problemáticas"),
        ("MEJORAR", "Trusted assets, vistas y contexto más preciso"),
    ]
    for index, (title, body) in enumerate(quality):
        angle_x = 0.72 + index * 2.5
        _box(
            slide,
            angle_x,
            2.05 if index % 2 == 0 else 3.15,
            2.15,
            1.52,
            fill=GRAY,
            line=(RED, ORANGE, TEAL, NAVY, RED)[index],
        )
        _textbox(
            slide,
            title,
            angle_x + 0.15,
            2.23 + (1.1 if index % 2 else 0),
            1.85,
            0.34,
            size=12,
            color=(RED, ORANGE, TEAL, NAVY, RED)[index],
            bold=True,
            align=PP_ALIGN.CENTER,
        )
        _textbox(
            slide,
            body,
            angle_x + 0.15,
            2.72 + (1.1 if index % 2 else 0),
            1.85,
            0.56,
            size=10,
            color=MID_GRAY,
            align=PP_ALIGN.CENTER,
            valign=MSO_ANCHOR.TOP,
        )
        if index < len(quality) - 1:
            _arrow(
                slide,
                (angle_x + 2.16, 2.83 + (1.1 if index % 2 else 0)),
                (angle_x + 2.48, 2.83 + (0 if index % 2 else 1.1)),
                color=LIGHT_LINE,
            )
    _textbox(
        slide,
        "Un Agent es un producto analítico vivo: se gestiona con evidencias, no solo con prompts.",
        1.2,
        5.65,
        10.9,
        0.52,
        size=18,
        color=TEAL,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    _source_note(slide, "Fuente: docs.databricks.com/aws/en/genie-agents/monitor/")
    _notes(
        slide,
        (
            "Posicionar benchmarks y monitorización como disciplina de producto, "
            "no como tarea puntual."
        ),
        "Aclaramos ahora qué significan las Apps dentro de Genie One.",
        "1:10",
    )

    # 16 · Apps
    slide = _dark_slide(
        presentation,
        "¿Qué son las Apps en el contexto de Genie One?",
        "Experiencias operativas personalizadas que viven sobre la plataforma",
    )
    _box(slide, 0.82, 1.55, 5.45, 4.7, fill=RGBColor(0x23, 0x30, 0x40), line=TEAL)
    _textbox(slide, "Databricks App", 1.2, 1.92, 4.7, 0.52, size=22, color=WHITE, bold=True)
    _textbox(
        slide,
        "Una aplicación web segura y serverless que combina:",
        1.2,
        2.62,
        4.65,
        0.5,
        size=14,
        color=LIGHT_LINE,
    )
    app_items = [
        "Interfaz adaptada al proceso de negocio",
        "Genie Agent como recurso con permiso CAN RUN",
        "SQL Warehouse y datos de Unity Catalog",
        "Workflows, modelos, secretos u otras Apps",
        "Identidad propia y permisos de mínimo privilegio",
    ]
    for index, item in enumerate(app_items):
        _textbox(slide, "✓", 1.25, 3.25 + index * 0.5, 0.3, 0.3, size=13, color=TEAL, bold=True)
        _textbox(slide, item, 1.68, 3.21 + index * 0.5, 4.05, 0.33, size=11, color=WHITE)
    _box(slide, 6.72, 1.55, 5.75, 4.7, fill=WHITE, line=ORANGE)
    _pill(slide, "EXTENSIÓN RECOMENDADA", 7.12, 1.92, 2.35, fill=ORANGE)
    _textbox(
        slide,
        "Renewable Action Center",
        7.12,
        2.58,
        4.75,
        0.5,
        size=21,
        color=NAVY,
        bold=True,
    )
    _textbox(
        slide,
        "Una futura App podría convertir el insight en acción:",
        7.12,
        3.22,
        4.75,
        0.4,
        size=13,
        color=MID_GRAY,
    )
    app_future = [
        "Consultar el Genie Agent",
        "Priorizar instalaciones de riesgo",
        "Registrar una recomendación operativa",
        "Lanzar un workflow o crear una orden",
    ]
    for index, item in enumerate(app_future):
        _textbox(
            slide,
            f"0{index + 1}",
            7.15,
            3.85 + index * 0.48,
            0.4,
            0.3,
            size=11,
            color=RED,
            bold=True,
        )
        _textbox(slide, item, 7.78, 3.81 + index * 0.48, 3.9, 0.34, size=11, bold=True)
    _textbox(
        slide,
        "No está desplegada en el demo actual; es el siguiente incremento lógico.",
        7.12,
        5.78,
        4.75,
        0.3,
        size=10,
        color=RED,
        bold=True,
    )
    _source_note(slide, "Fuente: docs.databricks.com/aws/en/dev-tools/databricks-apps/genie")
    _notes(
        slide,
        (
            "Distinguir una App de un Agent: la App implementa una experiencia "
            "y puede invocar al Agent."
        ),
        "Separamos claramente lo ya desplegado de las extensiones posibles.",
        "1:30",
    )

    # 17 · Now and next
    slide = _base_slide(
        presentation,
        "Capacidades actuales y evolución posible",
        "Un roadmap honesto: probado hoy frente a extensiones futuras",
    )
    _box(slide, 0.72, 1.42, 5.95, 4.95, fill=GRAY, line=TEAL)
    _pill(slide, "DESPLEGADO Y VALIDADO", 1.05, 1.72, 2.2, fill=TEAL)
    deployed = [
        "Tablas Delta y workflow serverless",
        "Semantic View + Metric View",
        "AI/BI Dashboard publicado",
        "Genie Agent con conversación real",
        "5 SQL de referencia + 5 benchmarks",
        "Resultado de benchmark: 100 % (5/5)",
        "Descubrimiento desde Genie One",
        "GitHub + Automation Bundle reproducible",
    ]
    for index, item in enumerate(deployed):
        _textbox(slide, "✓", 1.08, 2.45 + index * 0.43, 0.3, 0.28, size=12, color=TEAL, bold=True)
        _textbox(slide, item, 1.5, 2.41 + index * 0.43, 4.7, 0.3, size=10.5, bold=True)
    _box(slide, 6.95, 1.42, 5.65, 4.95, fill=WHITE, line=ORANGE)
    _pill(slide, "SIGUIENTES OPCIONES", 7.28, 1.72, 2.05, fill=ORANGE)
    optional = [
        ("App operativa", "Insight → workflow/acción"),
        ("Agent mode", "Investigaciones multiconsulta"),
        ("Archivos en Volumes", "Cruzar datos y documentos"),
        ("Conexiones externas", "Drive, Slack, Microsoft 365…"),
        ("Companion Genie", "Exploración desde el dashboard"),
        ("Integración externa", "Conversation API o embedding"),
    ]
    for index, (title, body) in enumerate(optional):
        y = 2.42 + index * 0.57
        _textbox(slide, title, 7.3, y, 1.75, 0.3, size=11, color=NAVY, bold=True)
        _textbox(slide, body, 9.15, y, 2.95, 0.3, size=10, color=MID_GRAY)
    _textbox(
        slide,
        "Recomendación: priorizar la App cuando sea necesario cerrar el ciclo de acción.",
        1.2,
        6.58,
        10.9,
        0.3,
        size=11,
        color=RED,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    _notes(
        slide,
        "Separar con transparencia lo demostrado de lo que forma parte del roadmap.",
        "Pasamos al modelo de acceso y entrega.",
        "1:15",
    )

    # 18 · Sharing
    slide = _dark_slide(
        presentation,
        "Acceso y entrega de los componentes",
        "La forma de acceso depende del rol, el entorno y los permisos",
    )
    sharing = [
        (
            "NEGOCIO",
            "Compartir en Databricks",
            "Dashboard publicado + Genie Agent + permisos\nConsumo desde Genie One",
            RED,
        ),
        (
            "EQUIPO TÉCNICO",
            "Entregar el repositorio",
            "Código, tests, definición del dashboard,\nconfiguración Genie y documentación",
            TEAL,
        ),
        (
            "OTRO WORKSPACE",
            "Desplegar el Automation Bundle",
            "El bundle reproduce jobs y dashboard;\nel script configura el Genie Agent",
            ORANGE,
        ),
        (
            "USUARIO EXTERNO",
            "PDF, vídeo o integración",
            "Embedding autenticado, App o Conversation API\nsegún identidad y permisos",
            NAVY,
        ),
    ]
    for index, (audience, title, body, color) in enumerate(sharing):
        row, column = divmod(index, 2)
        x, y = 0.82 + column * 6.12, 1.62 + row * 2.35
        _box(slide, x, y, 5.7, 1.92, fill=RGBColor(0x23, 0x30, 0x40), line=color)
        _pill(slide, audience, x + 0.25, y + 0.22, 1.45, fill=color)
        _textbox(slide, title, x + 1.9, y + 0.2, 3.35, 0.4, size=15, color=WHITE, bold=True)
        _textbox(
            slide,
            body,
            x + 0.25,
            y + 0.9,
            5.05,
            0.72,
            size=11,
            color=LIGHT_LINE,
            valign=MSO_ANCHOR.TOP,
        )
    _textbox(
        slide,
        "Regla práctica: negocio consume recursos publicados; tecnología recibe repo + bundle.",
        1.15,
        6.26,
        11.0,
        0.44,
        size=17,
        color=TEAL,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    _source_note(slide, "Fuente: docs.databricks.com/aws/en/dashboards/share/")
    _notes(
        slide,
        (
            "Responder explícitamente que el bundle no es el artefacto que recibe "
            "el usuario de negocio."
        ),
        "Cerramos con un recorrido de demo fácil de presentar.",
        "1:30",
    )

    # 19 · Client close
    slide = _dark_slide(
        presentation,
        "Una prueba de viabilidad que demuestra un patrón escalable",
        "Gobierno + BI + conversación + camino a la acción",
    )
    outcomes = [
        ("CONFIABLE", "Métricas y fuentes gobernadas", TEAL),
        ("CONVERSACIONAL", "Preguntas y follow-ups de negocio", ORANGE),
        ("REPRODUCIBLE", "Repo, tests y Automation Bundle", RED),
    ]
    for index, (title, body, color) in enumerate(outcomes):
        x = 0.9 + index * 4.12
        _box(slide, x, 1.72, 3.7, 2.08, fill=RGBColor(0x23, 0x30, 0x40), line=color)
        _textbox(slide, title, x + 0.2, 2.05, 3.3, 0.48, size=17, color=color, bold=True)
        _textbox(
            slide,
            body,
            x + 0.2,
            2.75,
            3.3,
            0.55,
            size=13,
            color=LIGHT_LINE,
            align=PP_ALIGN.CENTER,
        )
    _textbox(
        slide,
        "Siguiente decisión",
        1.1,
        4.65,
        2.35,
        0.42,
        size=13,
        color=RED,
        bold=True,
    )
    _textbox(
        slide,
        (
            "Validar el patrón con datos y usuarios reales; añadir una App "
            "solo si hay una acción operativa que cerrar."
        ),
        1.1,
        5.15,
        10.75,
        0.92,
        size=21,
        color=WHITE,
        bold=True,
        valign=MSO_ANCHOR.TOP,
    )
    _textbox(
        slide,
        "¿Qué decisión de negocio queremos acelerar primero?",
        1.1,
        6.35,
        8.4,
        0.42,
        size=15,
        color=TEAL,
        bold=True,
    )
    _notes(
        slide,
        "Cerrar pidiendo una decisión concreta y evitando prometer industrialización automática.",
        "Abrir preguntas.",
        "0:50",
    )

    # 20 · Technical appendix
    slide = _base_slide(
        presentation,
        "Anexo técnico · implementación del demo",
        "Job serverless, notebooks, lenguajes y controles que sostienen la experiencia",
    )
    _box(slide, 0.72, 1.36, 7.35, 5.42, fill=GRAY, line=NAVY)
    _pill(slide, "JOB SERVERLESS · 4 TAREAS", 1.03, 1.66, 2.45, fill=NAVY)
    notebook_tasks = [
        (
            "01",
            "01_setup_and_generate_data.py",
            "Crea el schema y publica tres tablas Delta sintéticas.",
            ORANGE,
        ),
        (
            "02",
            "02_transform_and_publish.py",
            "Publica KPI, Semantic View y Metric View si está disponible.",
            TEAL,
        ),
        (
            "03",
            "03_data_quality_checks.py",
            "Valida conteos, claves, rangos y contrato del dataset.",
            RED,
        ),
        (
            "04",
            "04_demo_validation.py",
            "Comprueba las vistas y emite la evidencia final del workflow.",
            NAVY,
        ),
    ]
    for index, (number, filename, purpose, color) in enumerate(notebook_tasks):
        y = 2.34 + index * 1.0
        _box(slide, 1.02, y, 6.73, 0.78, fill=WHITE, line=color)
        _textbox(slide, number, 1.2, y + 0.16, 0.52, 0.38, size=14, color=color, bold=True)
        _textbox(slide, filename, 1.82, y + 0.1, 2.68, 0.28, size=11.5, bold=True)
        _textbox(
            slide,
            purpose,
            4.58,
            y + 0.1,
            2.88,
            0.48,
            size=9.5,
            color=MID_GRAY,
            valign=MSO_ANCHOR.TOP,
        )
        if index < len(notebook_tasks) - 1:
            _arrow(slide, (4.35, y + 0.79), (4.35, y + 0.98), color=color)

    _box(slide, 8.35, 1.36, 4.27, 5.42, fill=NAVY, line=NAVY)
    _textbox(slide, "Stack real", 8.72, 1.74, 3.5, 0.42, size=19, color=WHITE, bold=True)
    stack_items = [
        "Python 3.11/3.12",
        "PySpark + Spark SQL",
        "Delta Lake + Unity Catalog",
        "Databricks SDK",
        "Automation Bundle · GitHub Actions",
    ]
    for index, item in enumerate(stack_items):
        _textbox(slide, "•", 8.75, 2.38 + index * 0.43, 0.25, 0.28, size=12, color=TEAL)
        _textbox(slide, item, 9.08, 2.34 + index * 0.43, 2.95, 0.3, size=10.5, color=WHITE)
    _textbox(slide, "Quality gates", 8.72, 4.74, 3.5, 0.36, size=16, color=TEAL, bold=True)
    quality_items = [
        "Ruff + Mypy",
        "Pytest: 29 passed · 92 % cobertura",
        "Smoke test + integración opt-in",
        "16 DQ + 8 SQL de smoke",
        "Benchmark Genie: 5/5",
    ]
    for index, item in enumerate(quality_items):
        _textbox(slide, "✓", 8.75, 5.22 + index * 0.29, 0.25, 0.22, size=10, color=TEAL)
        _textbox(slide, item, 9.08, 5.18 + index * 0.29, 3.05, 0.24, size=9, color=LIGHT_LINE)
    _notes(
        slide,
        (
            "Utilizar este anexo únicamente cuando haya audiencia técnica "
            "o preguntas de implementación."
        ),
        "La siguiente lámina explica cómo reproducir el proyecto en cada tipo de workspace.",
        "1:30",
    )

    # 21 · Reproduction by environment
    slide = _base_slide(
        presentation,
        "Cómo reproducirlo por entorno Databricks",
        "Mismo repositorio y bundle; cambian el target, los permisos y la autenticación",
    )
    _box(slide, 0.72, 1.42, 5.95, 4.78, fill=GRAY, line=TEAL)
    _pill(slide, "DATABRICKS FREE EDITION", 1.06, 1.75, 2.48, fill=TEAL)
    free_items = [
        ("Target", "dev · modo development"),
        ("Compute", "exclusivamente serverless"),
        ("Datos", "catalog workspace · schema aislado"),
        ("SQL", "warehouse existente · 2X-Small"),
        ("Auth", "OAuth U2M interactivo"),
        ("Uso", "prueba individual sujeta a fair-use"),
    ]
    for index, (label, body) in enumerate(free_items):
        y = 2.55 + index * 0.52
        _textbox(slide, label, 1.1, y, 1.0, 0.3, size=10.5, color=TEAL, bold=True)
        _textbox(slide, body, 2.15, y, 3.9, 0.3, size=10.5, color=NAVY)

    _box(slide, 6.95, 1.42, 5.65, 4.78, fill=WHITE, line=ORANGE)
    _pill(slide, "DATABRICKS ENTERPRISE", 7.29, 1.75, 2.38, fill=ORANGE)
    enterprise_items = [
        ("Target", "enterprise · modo production"),
        ("Compute", "serverless recomendado"),
        ("Datos", "catalog/schema corporativos autorizados"),
        ("SQL", "warehouse Pro/Serverless con CAN USE"),
        ("Auth", "U2M local · M2M en CI/CD"),
        ("Gobierno", "mínimo privilegio y permisos de bundle"),
    ]
    for index, (label, body) in enumerate(enterprise_items):
        y = 2.55 + index * 0.52
        _textbox(slide, label, 7.32, y, 1.0, 0.3, size=10.5, color=ORANGE, bold=True)
        _textbox(slide, body, 8.38, y, 3.68, 0.3, size=10.5, color=NAVY)

    _box(slide, 1.2, 6.23, 10.9, 0.48, fill=NAVY, line=NAVY)
    _textbox(
        slide,
        "1 · CLONAR + UV SYNC   →   2 · AUTH + VARIABLES   →   "
        "3 · VALIDATE + DEPLOY + RUN   →   4 · GENIE + SMOKE TEST",
        1.42,
        6.33,
        10.45,
        0.28,
        size=10.5,
        color=WHITE,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    _source_note(slide, "Detalle reproducible: README.md · docs/deployment_runbook.md")
    _notes(
        slide,
        "Aclarar que el bundle no crea infraestructura de cuenta ni un SQL warehouse.",
        "El runbook contiene los comandos exactos para ambos targets.",
        "1:20",
    )

    # 22 · Internal walkthrough, intentionally last
    slide = _base_slide(
        presentation,
        "Guía interna · recorrido recomendado",
        "Última diapositiva opcional; puede eliminarse antes de distribuir la versión final",
    )
    _pill(slide, "USO INTERNO", 10.9, 0.54, 1.45, fill=ORANGE)
    walkthrough = [
        ("01", "Contexto", "Explicar objetivo y datos sintéticos"),
        ("02", "KPIs", "Detectar desviación en Executive Overview"),
        ("03", "Diagnóstico", "Bajar a fiabilidad e instalaciones"),
        ("04", "Pregunta", "Consultar los tres riesgos en Genie"),
        ("05", "Evidencia", "Mostrar SQL, tabla y gráfico"),
        ("06", "Confianza", "Abrir Benchmark 5/5"),
        ("07", "Escala", "Enseñar Genie One y compartir"),
        ("08", "Acción", "Proponer App como siguiente paso"),
    ]
    for index, (number, title, body) in enumerate(walkthrough):
        row, column = divmod(index, 4)
        x, y = 0.58 + column * 3.18, 1.45 + row * 2.25
        _box(slide, x, y, 2.88, 1.78, fill=GRAY, line=(RED, NAVY, TEAL, ORANGE)[column])
        _textbox(slide, number, x + 0.16, y + 0.14, 0.52, 0.4, size=15, color=RED, bold=True)
        _textbox(slide, title, x + 0.78, y + 0.14, 1.82, 0.4, size=14, bold=True)
        _textbox(
            slide,
            body,
            x + 0.18,
            y + 0.77,
            2.46,
            0.65,
            size=10.5,
            color=MID_GRAY,
            valign=MSO_ANCHOR.TOP,
        )
    _textbox(
        slide,
        "ELIMINAR ESTA DIAPOSITIVA ANTES DE DISTRIBUIR LA VERSIÓN FINAL",
        1.35,
        6.23,
        10.65,
        0.45,
        size=14,
        color=RED,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    _notes(
        slide,
        (
            "Guía para el equipo presentador. Esta diapositiva se mantiene al final "
            "para que pueda eliminarse sin alterar la narrativa principal."
        ),
        "Fin del anexo interno.",
        "1:10",
    )

    return presentation


def validate_presentation(path: Path) -> dict[str, Any]:
    """Validate structure, notes, canvas bounds, and prohibited placeholders."""
    presentation = Presentation(path)
    failures: list[str] = []
    expected_slides = 22
    if len(presentation.slides) != expected_slides:
        failures.append(f"expected {expected_slides} slides, found {len(presentation.slides)}")
    prohibited = ("TODO", "Iberdrola")
    for slide_number, slide in enumerate(presentation.slides, start=1):
        notes = slide.notes_slide.notes_text_frame.text.strip()
        if not notes:
            failures.append(f"slide {slide_number} has no notes")
        for shape in slide.shapes:
            if shape.left < 0 or shape.top < 0:
                failures.append(f"slide {slide_number} has a negative shape position")
            if shape.left + shape.width > presentation.slide_width:
                failures.append(f"slide {slide_number} shape exceeds width")
            if shape.top + shape.height > presentation.slide_height:
                failures.append(f"slide {slide_number} shape exceeds height")
            if hasattr(shape, "text"):
                for token in prohibited:
                    if token.lower() in shape.text.lower():
                        failures.append(f"slide {slide_number} contains prohibited token {token}")
    if failures:
        raise ValueError("Presentation validation failed: " + "; ".join(failures))
    return {
        "slides": len(presentation.slides),
        "notes": len(presentation.slides),
        "canvas_failures": 0,
    }


def main() -> None:
    """Generate the presentation and local metric evidence."""
    metrics = _metric_summary()
    LOCAL_EVIDENCE.parent.mkdir(parents=True, exist_ok=True)
    serializable_metrics = {key: value for key, value in metrics.items() if key != "monthly"}
    LOCAL_EVIDENCE.write_text(
        json.dumps(serializable_metrics, indent=2, ensure_ascii=False), encoding="utf-8"
    )
    presentation = build_presentation(metrics)
    presentation.save(OUTPUT)
    validation = validate_presentation(OUTPUT)
    print(json.dumps({"output": str(OUTPUT), "validation": validation}, indent=2))


if __name__ == "__main__":
    main()
